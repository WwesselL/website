from flask import Flask, request, send_from_directory, redirect, url_for, jsonify, session, flash, get_flashed_messages, json, render_template, send_file
import os
from datetime import datetime, timedelta
import traceback
from functools import wraps
from werkzeug.utils import secure_filename
from flask_sqlalchemy import SQLAlchemy
import logging
import shutil
import json
import calendar
import secrets
import mimetypes
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import base64
import re
from collections import defaultdict
from openai import OpenAI

app = Flask(__name__, static_url_path='/static')
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
app.secret_key = 'jouw_geheime_sleutel_hier'  # Vervang dit door een veilige, willekeurige string
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///portfolio.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

common_styles = '''
body { 
    background-color: var(--bg-color); 
    color: var(--text-color);
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'Roboto', 'Helvetica', 'Arial', sans-serif; 
    margin: 0; 
    padding-top: 80px; /* Ruimte voor de header */
    transition: background-color 0.3s, color 0.3s;
}
:root {
    --bg-color: #ffffff;
    --text-color: #1d1d1f;
    --header-bg: #f5f5f7;
    --content-bg: #f5f5f7;
    --button-bg: #0077ed;  // Lichter blauw voor de lichte modus
    --button-text: #ffffff;
    --button-hover: #2b8ff7;  // Nog lichtere blauwe kleur voor hover
    --folder-bg: #e8e8ed;
    --folder-hover: #d2d2d7;
}
.dark-mode {
    --bg-color: #1d1d1f;
    --text-color: #f5f5f7;
    --header-bg: #2c2c2e;
    --content-bg: #2c2c2e;
    --button-bg: #0a84ff;  // Terug naar blauw voor dark mode
    --button-text: #ffffff;
    --button-hover: #409cff;  // Lichtere blauwe kleur voor hover in dark mode
    --folder-bg: #3a3a3c;
    --folder-hover: #4e4e50;
}
.header {
    position: fixed;
    top: 0;
    left: 0;
    right: 0;
    background-color: var(--header-bg);
    padding: 10px 0;
    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
    z-index: 1000;
}
.header-content {
    display: flex;
    justify-content: center;
    align-items: center;
    max-width: 1200px;
    margin: 0 auto;
}
.logo {
    margin-right: 20px;
    display: flex;
    align-items: center;
}
.logo img {
    height: 60px; /* Increased height */
    width: auto;
    margin-top: -10px; /* This will move the logo up */
    filter: var(--logo-filter);
}
.nav-buttons {
    display: flex;
    align-items: center; /* This ensures buttons stay vertically centered */
}
.nav-buttons button, .file-label, .delete-btn { 
    background-color: var(--button-bg); 
    color: var(--button-text); 
    border: none; 
    border-radius: 980px; 
    padding: 8px 16px; 
    font-size: 14px; 
    cursor: pointer; 
    transition: all 0.2s ease-in-out; 
    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1); 
    margin-right: 10px;
}
.nav-buttons button:hover, .file-label:hover { 
    background-color: var(--button-hover);
    transform: translateY(-2px); 
    box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2); 
}
.container { max-width: 1000px; margin: 0 auto; padding: 20px; }
h1 { font-size: 36px; margin-bottom: 20px; }
.content-container { 
    background-color: var(--content-bg); 
    padding: 30px; 
    border-radius: 18px; 
    margin-top: 20px; 
    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1); 
}
.hidden { display: none; }
a { color: #0066cc; text-decoration: none; }
a:hover { text-decoration: underline; }
.logo { position: fixed; top: 10px; left: 10px; font-size: 20px; font-weight: bold; z-index: 1000; }
.logo a { color: #0071e3; text-decoration: none; }
.logo a:hover { text-decoration: none; color: #0077ed; }
.logo img { height: 60px; width: auto; }
#darkModeToggle {
    background: none;
    border: none;
    color: var(--text-color);
    font-size: 24px;
    cursor: pointer;
    margin-left: 10px;
}
.file-structure {
    background-color: var(--content-bg);
    color: var(--text-color);
    padding: 20px;
    border-radius: 10px;
    display: flex;
}
.file-tree {
    list-style-type: none;
    padding-left: 0;
    flex: 0 0 250px;
}
.folder, .file {
    background-color: var(--folder-bg);
    color: var(--text-color);
    transition: background-color 0.3s, color 0.3s;
    margin-bottom: 5px;
    padding: 10px;
    border-radius: 5px;
}
.folder:hover, .file:hover {
    background-color: var(--folder-hover);
}
.folder span, .file a {
    color: var(--text-color);
    text-decoration: none;
    display: inline-block;
    margin-left: 5px;
}
.folder-item {
    position: relative;
}
.folder-content {
    position: absolute;
    left: calc(100% - 20px);  // Change this: Move 20px to the left
    top: 0;
    background-color: var(--content-bg);
    border-radius: 5px;
    padding: 5px;
    box-shadow: 0 2px 5px rgba(0, 0, 0, 0.1);
    z-index: 10;
    visibility: hidden;
    opacity: 0;
    transition: visibility 0s, opacity 0.3s linear;
    min-width: 200px;  // Add this: Ensure a minimum width for the subfolder content
}
.folder-content.visible {
    visibility: visible;
    opacity: 1;
}
.folder-item, .file-item {
    display: flex;
    align-items: center;
}
.folder-item button, .file-item button {
    margin-left: 10px;
}
.fas {
    width: 20px;
    text-align: center;
}
.folder {
    cursor: pointer;
}
'''

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'is_admin' not in session or not session['is_admin']:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

@app.route('/login', methods=['GET', 'POST'])
def login():
    error_message = ''
    if request.method == 'POST':
        if request.form['password'] == 'admin123':  # Vervang dit door een veilig wachtwoord
            session['is_admin'] = True
            return redirect(url_for('home'))
        else:
            error_message = 'Incorrect wachtwoord'
    
    is_admin = 'is_admin' in session and session['is_admin']
    header = generate_header(is_admin)
    
    return f'''
    <html>
        <head>
            <title>Admin Login - Portfolio Wessel</title>
            <style>
                {common_styles}
                .login-form {{
                    background-color: var(--content-bg);
                    padding: 30px;
                    border-radius: 18px;
                    margin-top: 60px;
                    max-width: 400px;
                    margin-left: auto;
                    margin-right: auto;
                }}
                .login-form input[type="password"] {{
                    width: 100%;
                    padding: 10px;
                    margin-bottom: 20px;
                    border: 1px solid #d2d2d7;
                    border-radius: 8px;
                    font-size: 16px;
                }}
                .login-form input[type="submit"] {{
                    background-color: var(--button-bg);
                    color: var(--button-text);
                    border: none;
                    border-radius: 980px;
                    padding: 10px 20px;
                    font-size: 16px;
                    cursor: pointer;
                    transition: background-color 0.2s;
                    width: 100%;
                }}
                .login-form input[type="submit"]:hover {{
                    background-color: var(--button-hover);
                }}
                .error-message {{
                    color: #ff3b30;
                    margin-bottom: 20px;
                }}
            </style>
        </head>
        <body>
            {header}
            <div class="container">
                <div class="login-form">
                    <h1>Admin Login</h1>
                    {f'<p class="error-message">{error_message}</p>' if error_message else ''}
                    <form method="post">
                        <input type="password" name="password" placeholder="Wachtwoord" required>
                        <input type="submit" value="Login">
                    </form>
                </div>
            </div>
            {generate_dark_mode_script()}
        </body>
    </html>
    '''

@app.route('/logout')
def logout():
    session.pop('is_admin', None)
    return redirect(url_for('home'))

def read_content(filename, month=None, year=None):
    if month and year:
        filename = f'logboek_{month}_{year}.txt'
    filepath = os.path.join('content', filename)
    if os.path.exists(filepath):
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()
    # Default content for September 2024
    if month == 9 and year == 2024:
        return '''
        <h3>Week 1</h3>
        <p>Introductie minor, kennismaking met medestudenten en docenten. Start met het opzetten van de portfolio website.</p>
        
        <h3>Week 2</h3>
        <p>Verdieping in webontwikkeling technieken. Verdere uitwerking van de portfolio structuur en design.</p>
        
        <h3>Week 3</h3>
        <p>Implementatie van basis functionaliteiten op de website, zoals navigatie en contentbeheer. Begin van het verzamelen van projectmateriaal.</p>
        '''
    return ''

def write_content(filename, content):
    filepath = os.path.join('content', filename)
    os.makedirs(os.path.dirname(filepath), exist_ok=True)
    with open(filepath, 'w', encoding='utf-8') as file:
        file.write(content)
    app.logger.info(f"Content written to {filepath}")

def generate_edit_button(is_admin):
    if is_admin:
        return '''
        <button id="editButton" class="edit-button" onclick="toggleEdit()">Bewerken</button>
        <style>
            .edit-button {
                position: fixed;
                bottom: 20px;
                left: 20px;
                z-index: 1000;
                background-color: var(--button-bg);
                color: var(--button-text);
                border: none;
                border-radius: 20px;
                padding: 10px 20px;
                font-size: 16px;
                cursor: pointer;
                transition: background-color 0.3s ease;
            }
            .edit-button:hover {
                background-color: var(--button-hover);
            }
        </style>
        '''
    return ''

def generate_edit_script():
    return '''
    <script>
        function toggleEdit() {
            const editableContent = document.querySelector('.editable-content');
            const editButton = document.getElementById('editButton');
            
            if (editButton.textContent === 'Bewerken') {
                const textarea = document.createElement('textarea');
                textarea.value = editableContent.innerHTML;
                textarea.style.width = '100%';
                textarea.style.height = '300px';
                editableContent.innerHTML = '';
                editableContent.appendChild(textarea);
                editButton.textContent = 'Opslaan';
            } else {
                const textarea = editableContent.querySelector('textarea');
                editableContent.innerHTML = textarea.value;
                editButton.textContent = 'Bewerken';
                saveChanges();
            }
        }

        function saveChanges() {
            const editableContent = document.querySelector('.editable-content');
            const currentDate = new Date();
            const data = {
                content: editableContent.innerHTML,
                month: currentDate.getMonth() + 1,
                year: currentDate.getFullYear()
            };

            fetch('/save_logboek', {
                method: 'POST',
                headers: {
                    'Content-Type': 'application/json',
                    'X-CSRFToken': getCsrfToken() // Zorg ervoor dat deze functie is gedefinieerd
                },
                body: JSON.stringify(data)
            })
            .then(response => {
                if (!response.ok) {
                    throw new Error('Server responded with status: ' + response.status);
                }
                return response.json();
            })
            .then(data => {
                if (data.success) {
                    alert('Wijzigingen opgeslagen');
                } else {
                    alert('Er is een fout opgetreden bij het opslaan: ' + data.message);
                }
            })
            .catch((error) => {
                console.error('Error:', error);
                alert('Er is een fout opgetreden bij het opslaan');
            });
        }
    </script>
    '''

@app.route('/save_logboek', methods=['POST'])
@admin_required
def save_logboek():
    try:
        data = request.json
        if not data or 'content' not in data or 'month' not in data or 'year' not in data:
            return jsonify({"success": False, "message": "Ongeldige invoer"}), 400

        content = data['content']
        month = data['month']
        year = data['year']

        filename = f'logboek_{month}_{year}.txt'
        write_content(filename, content)

        return jsonify({"success": True, "message": "Logboek succesvol opgeslagen"})
    except Exception as e:
        app.logger.error(f"Error in save_logboek: {str(e)}")
        return jsonify({"success": False, "message": f"Er is een fout opgetreden: {str(e)}"}), 500

def generate_header(is_admin):
    admin_button = '<button onclick="window.location.href=\'/logout\'">Logout</button>' if is_admin else '<button onclick="window.location.href=\'/login\'">Admin Login</button>'
    return f'''
    <header class="header">
        <div class="header-content">
            <div class="logo">
                <a href="/"><img src="{url_for('static', filename='logo.png')}" alt="Wessels Portfolio"></a>
            </div>
            <nav class="nav-buttons">
                <button onclick="window.location.href='/'">Homepagina</button>
                <button onclick="window.location.href='/bestanden'">Bestanden</button>
                <button onclick="window.location.href='/contact'">Contact</button>
                <button onclick="window.location.href='/logboek'">Logboek</button>
                {admin_button}
                <button id="darkModeToggle">🌓</button>
            </nav>
        </div>
    </header>
    '''

def generate_dark_mode_script():
    return '''
    <script>
        function toggleDarkMode() {
            document.body.classList.toggle('dark-mode');
            updateLogoFilter();
            localStorage.setItem('darkMode', document.body.classList.contains('dark-mode'));
        }

        function updateLogoFilter() {
            const logo = document.querySelector('.logo img');
            if (document.body.classList.contains('dark-mode')) {
                logo.style.filter = 'invert(1)';
            } else {
                logo.style.filter = 'none';
            }
        }

        document.addEventListener('DOMContentLoaded', () => {
            const darkModeToggle = document.getElementById('darkModeToggle');
            if (darkModeToggle) {
                darkModeToggle.addEventListener('click', toggleDarkMode);
            }
            
            if (localStorage.getItem('darkMode') === 'true') {
                document.body.classList.add('dark-mode');
                updateLogoFilter();
            }
        });
    </script>
    '''

@app.route('/')
def home():
    try:
        is_admin = 'is_admin' in session and session['is_admin']
        header = generate_header(is_admin)
        
        content = read_content('home.txt')
        if not content:
            content = '<h1>Welkom op mijn portfolio</h1><p>Dit is de homepagina van mijn portfolio website.</p>'

        return f'''
        <html>
            <head>
                <title>Homepagina - Portfolio Wessel</title>
                <style>
                    {common_styles}
                    .main-content {{
                        display: flex;
                        flex-direction: column;
                        margin-top: 60px;
                    }}
                    .top-content {{
                        display: flex;
                        justify-content: space-between;
                        margin-bottom: 40px;
                    }}
                    .welcome-section {{
                        width: 65%;
                        background-color: var(--content-bg);
                        padding: 40px;
                        border-radius: 18px;
                    }}
                    .welcome-section h1 {{
                        font-size: 48px;
                        line-height: 1.07143;
                        font-weight: 600;
                        letter-spacing: -.005em;
                        margin-bottom: 20px;
                    }}
                    .welcome-section p {{
                        font-size: 24px;
                        line-height: 1.10722;
                        font-weight: 400;
                        letter-spacing: .004em;
                        margin-bottom: 30px;
                    }}
                    .cta-button {{
                        background-color: var(--button-bg);
                        color: var(--button-text);
                        font-size: 18px;
                        line-height: 1.33337;
                        font-weight: 400;
                        letter-spacing: -.01em;
                        padding: 12px 24px;
                        border-radius: 980px;
                        display: inline-block;
                        text-decoration: none;
                        transition: all 0.2s ease-in-out;
                    }}
                    .cta-button:hover {{
                        background-color: var(--button-hover);
                        text-decoration: none;
                    }}
                    .contact-info {{
                        width: 30%;
                        background-color: var(--content-bg);
                        padding: 30px;
                        border-radius: 18px;
                        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                    }}
                    .contact-info h2 {{
                        font-size: 28px;
                        margin-bottom: 20px;
                    }}
                    .contact-info p {{
                        font-size: 16px;
                        line-height: 1.5;
                        margin-bottom: 10px;
                    }}
                    .features {{
                        display: flex;
                        justify-content: space-between;
                    }}
                    .feature {{
                        width: calc(33.33% - 20px);
                        text-align: center;
                        padding: 20px;
                        background-color: var(--bg-color);
                        border-radius: 18px;
                        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                    }}
                    .feature h2 {{
                        font-size: 24px;
                        line-height: 1.16667;
                        font-weight: 600;
                        letter-spacing: .009em;
                        margin-bottom: 15px;
                    }}
                    .feature p {{
                        font-size: 16px;
                        line-height: 1.47059;
                        font-weight: 400;
                        letter-spacing: -.022em;
                    }}
                </style>
            </head>
            <body>
                {header}
                <div class="container">
                    <div class="main-content">
                        <div class="top-content">
                            <div class="welcome-section">
                                <h1>Welkom op mijn portfolio</h1>
                                <p>Ontdek mijn werk en ervaring als student Minor I-DIMO</p>
                                <a href="/bestanden" class="cta-button">Bekijk mijn projecten</a>
                            </div>
                            <div class="contact-info">
                                <h2>Contactinformatie</h2>
                                <p><strong>Naam:</strong> Wessel Barendrecht</p>
                                <p><strong>E-mail:</strong> wessel.barendrecht@student.hu.nl</p>
                                <p><strong>Telefoon:</strong> +31 6 83183787</p>
                            </div>
                        </div>
                        <div class="features">
                            <div class="feature">
                                <h2>Innovatief</h2>
                                <p>Creatieve oplossingen voor complexe problemen</p>
                            </div>
                            <div class="feature">
                                <h2>Professioneel</h2>
                                <p>Hoogwaardige resultaten en efficiënte werkwijze</p>
                            </div>
                            <div class="feature">
                                <h2>Toekomstgericht</h2>
                                <p>Altijd op zoek naar de nieuwste technologieën</p>
                            </div>
                        </div>
                    </div>
                </div>
                {generate_dark_mode_script()}
            </body>
        </html>
        '''
    except Exception as e:
        logging.error(f"Error in home route: {str(e)}")
        return f"Er is een fout opgetreden: {str(e)}", 500

def get_week_of_month(dt):
    first_day = dt.replace(day=1)
    dom = dt.day
    adjusted_dom = dom + first_day.weekday()
    return (adjusted_dom - 1) // 7 + 1

def get_upload_path(filename):
    now = datetime.now()
    year_month = now.strftime("%Y-%m")
    week = get_week_of_month(now)
    path = os.path.join(app.config['UPLOAD_FOLDER'], year_month, f"week_{week}")
    os.makedirs(path, exist_ok=True)
    return os.path.join(path, filename)

# Functie om de uploaddatum van een bestand op te slaan
def save_upload_date(file_path):
    upload_info_path = file_path + '.upload_info'
    with open(upload_info_path, 'w') as f:
        f.write(str(datetime.now().timestamp()))

# Functie om de uploaddatum van een bestand te lezen
def get_upload_date(file_path):
    upload_info_path = file_path + '.upload_info'
    if os.path.exists(upload_info_path):
        with open(upload_info_path, 'r') as f:
            return float(f.read().strip())
    return os.path.getctime(file_path)  # Fallback to creation time if upload info doesn't exist

def get_folder_creation_date(path):
    return os.path.getctime(path)

def get_folder_structure(path):
    structure = []
    for item in os.listdir(path):
        item_path = os.path.join(path, item)
        if os.path.isdir(item_path):
            structure.append({
                'name': item,
                'type': 'folder',
                'items': get_folder_structure(item_path),
                'creation_date': get_folder_creation_date(item_path)
            })
        else:
            structure.append({
                'name': item,
                'type': 'file',
                'upload_date': get_upload_date(item_path)
            })
    
    # Sorteer de structuur op basis van de opgeslagen volgorde
    if os.path.exists('folder_order.json'):
        with open('folder_order.json', 'r') as f:
            folder_order = json.load(f)
        structure.sort(key=lambda x: folder_order.index(x['name']) if x['name'] in folder_order else float('inf'))
    else:
        # Als er geen opgeslagen volgorde is, sorteer dan op datum (nieuwste eerst)
        structure.sort(key=lambda x: x.get('creation_date', 0) if x['type'] == 'folder' else x.get('upload_date', 0), reverse=True)
    
    return structure

def generate_structure_html(structure, base_path='', is_admin=False):
    html = '<ul class="file-tree">'
    for item in structure:
        if item['type'] == 'folder':
            folder_path = os.path.join(base_path, item['name'])
            html += f'''
                <li class="folder-item">
                    <span class="folder" onclick="toggleFolder(this)" draggable="true" data-path="{folder_path}">
                        <i class="fas fa-folder"></i> {item['name']}
                    </span>
                    {f'<button onclick="addFolder(\'{folder_path}\')">Nieuwe submap</button>' if is_admin else ''}
                    {f'<button class="delete-btn" onclick="deleteFolder(\'{folder_path}\')"><i class="fas fa-trash"></i> Verwijder map</button>' if is_admin else ''}
                    <div class="folder-content">
                        {generate_structure_html(item.get('items', []), folder_path, is_admin)}
                    </div>
                </li>
            '''
        else:
            file_path = os.path.join(base_path, item['name'])
            html += f'''
                <li class="file-item">
                    <span class="file">
                        <i class="fas fa-file"></i> 
                        <a href="#" onclick="previewFile('{file_path}', '{item['name']}'); return false;">{item['name']}</a>
                    </span>
                    {f'<button class="delete-btn" onclick="deleteFile(\'{file_path}\')"><i class="fas fa-trash"></i> Verwijderen</button>' if is_admin else ''}
                </li>
            '''
    html += '</ul>'
    return html

def generate_csrf_token():
    if 'csrf_token' not in session:
        session['csrf_token'] = secrets.token_hex(16)
    return session['csrf_token']

app.jinja_env.globals['csrf_token'] = generate_csrf_token

def verify_csrf_token():
    token = request.headers.get('X-CSRFToken') or request.form.get('csrf_token')
    return token and token == session.get('csrf_token')

@app.route('/bestanden', methods=['GET', 'POST'])
def bestanden():
    logging.debug("Entering bestanden route")
    is_admin = 'is_admin' in session and session['is_admin']
    header = generate_header(is_admin)
    
    if request.method == 'POST' and is_admin:
        action = request.json.get('action')
        if action == 'add_folder':
            folder_name = request.json.get('folder_name')
            parent_path = request.json.get('parent_path', '')
            try:
                new_folder_path = os.path.join(app.config['UPLOAD_FOLDER'], parent_path, folder_name)
                os.makedirs(new_folder_path, exist_ok=True)
                
                # Update the file structure data
                update_file_structure()
                
                return jsonify({
                    "success": True,
                    "message": "Map succesvol toegevoegd!",
                    "new_folder": {"name": folder_name, "type": "folder", "items": []}
                })
            except Exception as e:
                app.logger.error(f"Fout bij het maken van map: {str(e)}")
                return jsonify({
                    "success": False,
                    "message": f"Fout bij het maken van map: {str(e)}"
                }), 500

    return f'''
    <html>
        <head>
            <title>Bestanden - Portfolio Wessel</title>
            <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/5.15.3/css/all.min.css">
            <meta name="csrf-token" content="{generate_csrf_token()}">
            <style>
                {common_styles}
                /* Voeg hier de dark mode variabelen toe */
                :root {{
                    --bg-color: #ffffff;
                    --text-color: #1d1d1f;
                    --header-bg: #f5f5f7;
                    --content-bg: #f5f5f7;
                    --button-bg: #0077ed;  // Lichter blauw voor de lichte modus
                    --button-text: #ffffff;
                    --button-hover: #2b8ff7;  // Nog lichtere blauwe kleur voor hover
                    --folder-bg: #e8e8ed;
                    --folder-hover: #d2d2d7;
                }}
                .dark-mode {{
                    --bg-color: #1d1d1f;
                    --text-color: #f5f5f7;
                    --header-bg: #2c2c2e;
                    --content-bg: #2c2c2e;
                    --button-bg: #0a84ff;  // Terug naar blauw voor dark mode
                    --button-text: #ffffff;
                    --button-hover: #409cff;  // Lichtere blauwe kleur voor hover in dark mode
                    --folder-bg: #3a3a3c;
                    --folder-hover: #4e4e50;
                }}
                body {{
                    background-color: var(--bg-color);
                    color: var(--text-color);
                }}
                .file-structure {{
                    background-color: var(--content-bg);
                    color: var(--text-color);
                }}
                .folder, .file {{
                    background-color: var(--folder-bg);
                    color: var(--text-color);
                    transition: background-color 0.3s, color 0.3s;
                }}
                .folder:hover, .file:hover {{
                    background-color: var(--folder-hover);
                }}
                .folder span, .file a {{
                    color: var(--text-color);
                }}
                .folder-content {{
                    position: absolute;
                    left: calc(100% - 20px);  // Change this: Move 20px to the left
                    top: 0;
                    background-color: var(--content-bg);
                    border-radius: 5px;
                    padding: 5px;
                    box-shadow: 0 2px 5px rgba(0, 0, 0, 0.1);
                    z-index: 10;
                    visibility: hidden;
                    opacity: 0;
                    transition: visibility 0s, opacity 0.3s linear;
                    min-width: 200px;
                }}
                .folder-content.visible {{
                    visibility: visible;
                    opacity: 1;
                }}
                .nav-buttons, .admin-controls {{
                    display: flex;
                    justify-content: flex-start;
                    margin-bottom: 15px;
                }}
                .nav-buttons button, .admin-controls button, .file-label, .delete-btn {{
                    margin-right: 10px;
                    background-color: var(--button-bg); 
                    color: var(--button-text); 
                    border: none; 
                    border-radius: 980px; 
                    padding: 8px 16px; 
                    font-size: 14px; 
                    cursor: pointer; 
                    transition: all 0.2s ease-in-out; 
                    box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
                }}
                .nav-buttons button:hover, .admin-controls button:hover, .file-label:hover, .delete-btn:hover {{
                    background-color: var(--button-hover);
                    transform: translateY(-2px); 
                    box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2); 
                }}
                .delete-btn {{
                    background-color: #ff3b30;
                }}
                .delete-btn:hover {{
                    background-color: #ff453a;
                }}
                .file-label {{
                    display: inline-block;
                }}
                .file-input {{
                    display: none;
                }}
                .modal {{
                    display: none;
                    position: fixed;
                    z-index: 1000;
                    left: 0;
                    top: 0;
                    width: 100%;
                    height: 100%;
                    overflow: auto;
                    background-color: rgba(0,0,0,0.4);
                }}
                .modal-content {{
                    background-color: #fefefe;
                    margin: 15% auto;
                    padding: 20px;
                    border: 1px solid #888;
                    width: 80%;
                    max-width: 700px;
                }}
                .close {{
                    color: #aaa;
                    float: right;
                    font-size: 28px;
                    font-weight: bold;
                    cursor: pointer;
                }}
                .close:hover,
                .close:focus {{
                    color: black;
                    text-decoration: none;
                    cursor: pointer;
                }}
                #preview-content {{
                    max-width: 100%;
                    max-height: 400px;
                    overflow: auto;
                }}
                #preview-content img {{
                    max-width: 100%;
                    height: auto;
                }}
                .preview-docx, .preview-pptx {{
                    max-width: 100%;
                    overflow-x: auto;
                }}
                .preview-docx p, .preview-pptx p {{
                    margin-bottom: 10px;
                }}
                .preview-pptx h2 {{
                    margin-top: 20px;
                    margin-bottom: 10px;
                }}
                .pptx-preview .slide {{
                    border: 1px solid #ddd;
                    margin-bottom: 20px;
                    padding: 10px;
                }}
                .pptx-preview h2 {{
                    color: #333;
                    border-bottom: 1px solid #ddd;
                }}
                .pptx-preview h3 {{
                    color: #666;
                }}
                .pptx-preview table {{
                    border-collapse: collapse;
                    width: 100%;
                }}
                .pptx-preview td {{
                    padding: 5px;
                    border: 1px solid #ddd;
                }}
                .folder {{
                    cursor: move;
                }}
                .folder.dragging {{
                    opacity: 0.5;
                }}
                .folder-placeholder {{
                    border: 2px dashed #0071e3;
                    margin: 5px 0;
                    height: 30px;
                    border-radius: 8px;
                }}
            </style>
        </head>
        <body>
            {header}
            <div class="container">
                <h1>Bestandsbeheer</h1>
                {f'''
                <div class="admin-controls">
                    <button onclick="addFolder('')">Nieuwe map toevoegen</button>
                </div>
                ''' if is_admin else ''}
                <div class="file-structure" id="file-structure">
                    <!-- Hier komt de bestandsstructuur -->
                </div>
            </div>
            
            <div id="preview-modal" class="modal">
                <div class="modal-content">
                    <span class="close">&times;</span>
                    <h2 id="preview-title"></h2>
                    <div id="preview-content"></div>
                    <a id="download-link" href="#" download>Download bestand</a>
                </div>
            </div>
            
            <script>
                console.log("Starting JavaScript execution");
                const csrfToken = document.querySelector('meta[name="csrf-token"]').getAttribute('content');
                const isAdmin = {str(is_admin).lower()};
                let fileStructureData = []; // Initialize as an empty array

                function loadFileStructure() {{
                    console.log("Loading file structure");
                    fetch('/get_file_structure')
                        .then(response => response.json())
                        .then(data => {{
                            console.log("Received file structure data:", data);
                            fileStructureData = data;
                            updateFileStructureDisplay();
                        }})
                        .catch(error => {{
                            console.error('Error loading file structure:', error);
                            fileStructureData = [];
                            updateFileStructureDisplay();
                        }});
                }}

                function updateFileStructureDisplay() {{
                    console.log("Updating file structure display");
                    const fileStructure = document.getElementById('file-structure');
                    fileStructure.innerHTML = generateFileStructureHTML(fileStructureData);
                }}

                function generateFileStructureHTML(structure, currentPath = '') {{
                    if (!Array.isArray(structure)) return '';
                    let html = '<ul class="file-tree">';
                    for (const item of structure) {{
                        const itemPath = currentPath ? `${{currentPath}}/${{item.name}}` : item.name;
                        if (item.type === 'folder') {{
                            html += `
                                <li class="folder-item">
                                    <span class="folder" onclick="toggleFolder(this)" draggable="true" data-path="${{itemPath}}">
                                        <i class="fas fa-folder"></i> ${{item.name}}
                                    </span>
                                    ${{isAdmin ? `
                                        <button onclick="addFolder('${{itemPath}}')">Nieuwe submap</button>
                                        <button class="delete-btn" onclick="deleteFolder('${{itemPath}}')"><i class="fas fa-trash"></i></button>
                                        <label class="file-label">
                                            <i class="fas fa-upload"></i> Upload
                                            <input type="file" class="file-input" onchange="uploadFile(this, '${{itemPath}}')" />
                                        </label>
                                    ` : ''}}
                                    <div class="folder-content">
                                        ${{generateFileStructureHTML(item.items || [], itemPath)}}
                                    </div>
                                </li>`;
                        }} else {{
                            html += `
                                <li class="file-item">
                                    <span class="file">
                                        <i class="fas fa-file"></i>
                                        <a href="#" onclick="previewFile('${{itemPath}}', '${{item.name}}'); return false;">${{item.name}}</a>
                                    </span>
                                    ${{isAdmin ? `
                                        <button class="delete-btn" onclick="deleteFile('${{itemPath}}')"><i class="fas fa-trash"></i></button>
                                    ` : ''}}
                                </li>`;
                        }}
                    }}
                    html += '</ul>';
                    return html;
                }}

                function toggleFolder(element) {{
                    const folderContent = element.parentElement.querySelector('.folder-content');
                    if (folderContent) {{
                        folderContent.classList.toggle('visible');
                        const isVisible = folderContent.classList.contains('visible');
                        element.querySelector('i').className = isVisible ? 'fas fa-folder-open' : 'fas fa-folder';
                        
                        // Sluit andere open mappen op hetzelfde niveau
                        const parentUl = element.closest('ul');
                        parentUl.querySelectorAll(':scope > li > .folder-content.visible').forEach(content => {{
                            if (content !== folderContent) {{
                                content.classList.remove('visible');
                                content.previousElementSibling.querySelector('i').className = 'fas fa-folder';
                            }}
                        }});
                    }}
                }}

                function addFolder(parentPath) {{
                    const folderName = prompt('Voer de naam van de nieuwe map in:');
                    if (folderName) {{
                        fetch('/bestanden', {{
                            method: 'POST',
                            headers: {{
                                'Content-Type': 'application/json',
                                'X-CSRFToken': csrfToken
                            }},
                            body: JSON.stringify({{
                                action: 'add_folder',
                                folder_name: folderName,
                                parent_path: parentPath
                            }})
                        }})
                        .then(response => {{
                            if (!response.ok) {{
                                throw new Error('Server responded with an error');
                            }}
                            return response.json();
                        }})
                        .then(data => {{
                            if (data.success) {{
                                addFolderToStructure(parentPath, data.new_folder);
                                updateFileStructureDisplay();
                                alert(data.message);
                            }} else {{
                                throw new Error(data.message || 'Er is een onbekende fout opgetreden');
                            }}
                        }})
                        .catch(error => {{
                            console.error('Error:', error);
                            alert('Er is een fout opgetreden bij het aanmaken van de map: ' + error.message);
                        }});
                    }}
                }}

                function addFolderToStructure(parentPath, newFolder) {{
                    if (!Array.isArray(fileStructureData)) {{
                        console.error('fileStructureData is not an array');
                        fileStructureData = [];
                    }}
                    if (parentPath === '') {{
                        fileStructureData.push(newFolder);
                    }} else {{
                        const pathParts = parentPath.split('/');
                        let currentLevel = fileStructureData;
                        for (const part of pathParts) {{
                            let folder = currentLevel.find(item => item.type === 'folder' && item.name === part);
                            if (!folder) {{
                                folder = {{ name: part, type: 'folder', items: [] }};
                                currentLevel.push(folder);
                            }}
                            if (!Array.isArray(folder.items)) {{
                                folder.items = [];
                            }}
                            currentLevel = folder.items;
                        }}
                        currentLevel.push(newFolder);
                    }}
                }}

                function deleteFolder(folderPath) {{
                    if (confirm(`Weet je zeker dat je de map "${{folderPath}}" wilt verwijderen? Dit kan niet ongedaan worden gemaakt.`)) {{
                        fetch('/verwijder_map', {{
                            method: 'POST',
                            headers: {{
                                'Content-Type': 'application/json',
                                'X-CSRFToken': csrfToken
                            }},
                            body: JSON.stringify({{ folder_path: folderPath }})
                        }})
                        .then(response => response.json())
                        .then(data => {{
                            if (data.success) {{
                                alert(data.message);
                                loadFileStructure();  // Reload the file structure
                            }} else {{
                                alert('Fout bij het verwijderen van de map: ' + data.message);
                            }}
                        }})
                        .catch(error => {{
                            console.error('Error:', error);
                            alert('Er is een fout opgetreden bij het verwijderen van de map');
                        }});
                    }}
                }}

                function uploadFile(input, folderPath) {{
                    const file = input.files[0];
                    if (file) {{
                        const formData = new FormData();
                        formData.append('file', file);
                        formData.append('folder_path', folderPath);

                        fetch('/upload_bestand', {{
                            method: 'POST',
                            headers: {{
                                'X-CSRFToken': csrfToken
                            }},
                            body: formData
                        }})
                        .then(response => response.json())
                        .then(data => {{
                            if (data.success) {{
                                alert(data.message);
                                loadFileStructure();  // Reload the file structure
                            }} else {{
                                alert('Fout bij het uploaden van het bestand: ' + data.message);
                            }}
                        }})
                        .catch(error => {{
                            console.error('Error:', error);
                            alert('Er is een fout opgetreden bij het uploaden van het bestand');
                        }});
                    }}
                }}

                function deleteFile(filePath) {{
                    if (confirm(`Weet je zeker dat je het bestand "${{filePath}}" wilt verwijderen? Dit kan niet ongedaan worden gemaakt.`)) {{
                        fetch('/verwijder_bestand', {{
                            method: 'POST',
                            headers: {{
                                'Content-Type': 'application/json',
                                'X-CSRFToken': csrfToken
                            }},
                            body: JSON.stringify({{ file_path: filePath }})
                        }})
                        .then(response => response.json())
                        .then(data => {{
                            if (data.success) {{
                                alert(data.message);
                                loadFileStructure();  // Reload the file structure
                            }} else {{
                                alert('Fout bij het verwijderen van het bestand: ' + data.message);
                            }}
                        }})
                        .catch(error => {{
                            console.error('Error:', error);
                            alert('Er is een fout opgetreden bij het verwijderen van het bestand');
                        }});
                    }}
                }}

                function previewFile(filePath, fileName) {{
                    const modal = document.getElementById('preview-modal');
                    const previewTitle = document.getElementById('preview-title');
                    const previewContent = document.getElementById('preview-content');
                    const downloadLink = document.getElementById('download-link');
                    
                    previewTitle.textContent = fileName;
                    downloadLink.href = `/download/${{filePath}}`;
                    downloadLink.download = fileName;
                    
                    fetch(`/preview/${{filePath}}`)
                        .then(response => {{
                            if (!response.ok) {{
                                return response.json().then(err => {{ throw err; }});
                            }}
                            return response.text();
                        }})
                        .then(data => {{
                            if (data.startsWith('<img')) {{
                                previewContent.innerHTML = data;
                            }} else if (fileName.endsWith('.docx')) {{
                                previewContent.innerHTML = `<div class="preview-docx">${{data}}</div>`;
                            }} else if (fileName.endsWith('.pptx')) {{
                                previewContent.innerHTML = `<div class="preview-pptx">${{data}}</div>`;
                            }} else {{
                                previewContent.innerHTML = `<pre>${{data}}</pre>`;
                            }}
                            modal.style.display = 'block';
                        }})
                        .catch(error => {{
                            console.error('Error:', error);
                            previewContent.textContent = `Fout bij het laden van de preview: ${{error.message || 'Onbekende fout'}}`;
                            modal.style.display = 'block';
                        }});
                }}

                // Close the modal when clicking on <span> (x)
                document.querySelector('.close').onclick = function() {{
                    document.getElementById('preview-modal').style.display = 'none';
                }}

                // Close the modal when clicking outside of it
                window.onclick = function(event) {{
                    const modal = document.getElementById('preview-modal');
                    if (event.target == modal) {{
                        modal.style.display = 'none';
                    }}
                }}

                let draggedItem = null;

                function initDragAndDrop() {{
                    const fileStructure = document.querySelector('.file-structure');
                    fileStructure.addEventListener('dragstart', dragStart);
                    fileStructure.addEventListener('dragover', dragOver);
                    fileStructure.addEventListener('drop', drop);
                    fileStructure.addEventListener('dragend', dragEnd);
                }}

                function dragStart(e) {{
                    if (e.target.classList.contains('folder')) {{
                        draggedItem = e.target;
                        setTimeout(() => e.target.style.opacity = '0.5', 0);
                    }}
                }}

                function dragOver(e) {{
                    e.preventDefault();
                    if (e.target.classList.contains('folder') && e.target !== draggedItem) {{
                        const rect = e.target.getBoundingClientRect();
                        const midY = rect.top + rect.height / 2;
                        if (e.clientY < midY) {{
                            e.target.style.borderTop = '2px solid #0071e3';
                            e.target.style.borderBottom = '';
                        }} else {{
                            e.target.style.borderBottom = '2px solid #0071e3';
                            e.target.style.borderTop = '';
                        }}
                    }}
                }}

                function drop(e) {{
                    e.preventDefault();
                    if (e.target.classList.contains('folder') && e.target !== draggedItem) {{
                        const rect = e.target.getBoundingClientRect();
                        const midY = rect.top + rect.height / 2;
                        const parentUl = e.target.closest('ul');
                        if (e.clientY < midY) {{
                            parentUl.insertBefore(draggedItem.parentElement, e.target.parentElement);
                        }} else {{
                            parentUl.insertBefore(draggedItem.parentElement, e.target.parentElement.nextSibling);
                        }}
                        updateFolderOrder();
                    }}
                }}

                function dragEnd() {{
                    draggedItem.style.opacity = '1';
                    document.querySelectorAll('.folder').forEach(folder => {{
                        folder.style.borderTop = '';
                        folder.style.borderBottom = '';
                    }});
                    draggedItem = null;
                }}

                function updateFolderOrder() {{
                    const folderOrder = Array.from(document.querySelectorAll('.folder'))
                        .map(folder => folder.dataset.path);
                    
                    fetch('/update_folder_order', {{
                        method: 'POST',
                        headers: {{
                            'Content-Type': 'application/json',
                            'X-CSRFToken': csrfToken
                        }},
                        body: JSON.stringify({{ folder_order: folderOrder }})
                    }})
                    .then(response => response.json())
                    .then(data => {{
                        if (data.success) {{
                            console.log('Folder order updated successfully');
                        }} else {{
                            console.error('Failed to update folder order:', data.message);
                        }}
                    }})
                    .catch(error => {{
                        console.error('Error updating folder order:', error);
                    }});
                }}

                function toggleDarkMode() {{
                    console.log('toggleDarkMode called');
                    document.body.classList.toggle('dark-mode');
                    updateLogoFilter();
                    localStorage.setItem('darkMode', document.body.classList.contains('dark-mode'));
                    console.log('Dark mode is now:', document.body.classList.contains('dark-mode'));
                }}

                function updateLogoFilter() {{
                    console.log('updateLogoFilter called');
                    const logo = document.querySelector('.logo img');
                    if (logo) {{
                        logo.style.filter = document.body.classList.contains('dark-mode') ? 'invert(1)' : 'none';
                        console.log('Logo filter updated');
                    }} else {{
                        console.log('Logo element not found');
                    }}
                }}

                function initDarkMode() {{
                    console.log('initDarkMode called');
                    const darkModeToggle = document.getElementById('darkModeToggle');
                    if (darkModeToggle) {{
                        console.log('Dark mode toggle button found');
                        darkModeToggle.addEventListener('click', toggleDarkMode);
                    }} else {{
                        console.log('Dark mode toggle button not found');
                    }}

                    if (localStorage.getItem('darkMode') === 'true') {{
                        console.log('Dark mode should be enabled based on localStorage');
                        document.body.classList.add('dark-mode');
                        updateLogoFilter();
                    }}
                }}

                document.addEventListener('DOMContentLoaded', () => {{
                    console.log("DOM content loaded");
                    loadFileStructure();
                    initDragAndDrop();
                    initDarkMode();
                }});

                console.log('Script ended on bestanden page');
            </script>
        </body>
    </html>
    '''

@app.route('/get_file_structure')
def get_file_structure():
    logging.debug("Entering get_file_structure route")
    logging.debug(f"Current file structure: {file_structure}")
    return jsonify(file_structure)

def update_file_structure():
    global file_structure
    file_structure = get_folder_structure(app.config['UPLOAD_FOLDER'])

def generate_folder_options(structure, parent_path=''):
    options = ''
    for i, folder in enumerate(structure):
        current_path = f"{parent_path}/{folder['name']}" if parent_path else folder['name']
        options += f'<option value="{current_path}">{current_path}</option>'
        options += generate_folder_options(folder.get('items', []), current_path)
    return options

@app.route('/verwijder_bestand', methods=['POST'])
@admin_required
def verwijder_bestand():
    if not verify_csrf_token():
        return jsonify({"success": False, "message": "Ongeldige CSRF-token"}), 403
    
    file_path = request.json.get('file_path')
    if file_path:
        full_path = os.path.join(app.config['UPLOAD_FOLDER'], file_path)
        if os.path.exists(full_path):
            try:
                os.remove(full_path)
                # Verwijder ook het .upload_info bestand als het bestaat
                upload_info_path = full_path + '.upload_info'
                if os.path.exists(upload_info_path):
                    os.remove(upload_info_path)
                update_file_structure()  # Update the file structure after deletion
                return jsonify({"success": True, "message": "Bestand succesvol verwijderd", "file_path": file_path})
            except Exception as e:
                app.logger.error(f"Fout bij verwijderen bestand: {str(e)}")
                return jsonify({"success": False, "message": f"Fout bij verwijderen: {str(e)}"}), 500
        else:
            return jsonify({"success": False, "message": "Bestand niet gevonden"}), 404
    else:
        return jsonify({"success": False, "message": "Geen bestandspad opgegeven"}), 400

@app.route('/verwijder_map', methods=['POST'])
@admin_required
def verwijder_map():
    if not verify_csrf_token():
        return jsonify({"success": False, "message": "Ongeldige CSRF-token"}), 403
    
    folder_path = request.json.get('folder_path')
    if folder_path:
        full_path = os.path.join(app.config['UPLOAD_FOLDER'], folder_path)
        if os.path.exists(full_path) and os.path.isdir(full_path):
            try:
                shutil.rmtree(full_path)
                update_file_structure()  # Update the file structure after deletion
                return jsonify({"success": True, "message": "Map succesvol verwijderd"})
            except Exception as e:
                return jsonify({"success": False, "message": str(e)})
        else:
            return jsonify({"success": False, "message": "Map niet gevonden"})
    else:
        return jsonify({"success": False, "message": "Geen mappad opgegeven"})

@app.route('/upload_bestand', methods=['POST'])
@admin_required
def upload_bestand():
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "Geen bestand geüpload"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"success": False, "message": "Geen bestand geselecteerd"}), 400
    if file:
        filename = secure_filename(file.filename)
        folder_path = request.form.get('folder_path', '')
        upload_path = os.path.join(app.config['UPLOAD_FOLDER'], folder_path, filename)
        os.makedirs(os.path.dirname(upload_path), exist_ok=True)
        file.save(upload_path)
        save_upload_date(upload_path)
        update_file_structure()
        return jsonify({"success": True, "message": "Bestand succesvol geüpload", "file_path": os.path.join(folder_path, filename)})

@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

def docx_to_html(file_path):
    doc = docx.Document(file_path)
    html = "<div>"
    for para in doc.paragraphs:
        html += f"<p>{para.text}</p>"
    html += "</div>"
    return html

def pptx_to_html(file_path):
    prs = Presentation(file_path)
    html = "<div class='pptx-preview'>"
    
    for i, slide in enumerate(prs.slides, start=1):
        html += f"<div class='slide'><h2>Slide {i}</h2>"
        
        if slide.shapes.title:
            html += f"<h3>{slide.shapes.title.text}</h3>"
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, 'text') and shape.text.strip():
                    html += f"<p>{shape.text}</p>"
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                html += "<table border='1'>"
                for row in shape.table.rows:
                    html += "<tr>"
                    for cell in row.cells:
                        html += f"<td>{cell.text}</td>"
                    html += "</tr>"
                html += "</table>"
        
        html += "</div>"
    
    html += "</div>"
    return html

@app.route('/preview/<path:filename>')
def preview_file(filename):
    try:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if not os.path.exists(file_path):
            return jsonify({"success": False, "message": "Bestand niet gevonden"}), 404

        mime_type, _ = mimetypes.guess_type(file_path)
        
        if mime_type and mime_type.startswith('image/'):
            return f'<img src="{url_for("uploaded_file", filename=filename)}" alt="{filename}">'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return docx_to_html(file_path)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return pptx_to_html(file_path)
        elif mime_type and mime_type.startswith('text/'):
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            return content
        else:
            return "Preview niet beschikbaar voor dit bestandstype."
    except Exception as e:
        logging.error(f"Error in preview_file: {str(e)}")
        return jsonify({"success": False, "message": f"Er is een fout opgetreden: {str(e)}"}), 500

@app.route('/download/<path:filename>')
def download_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=True)

@app.errorhandler(Exception)
def handle_exception(e):
    app.logger.error(f"Onverwachte fout: {str(e)}")
    return jsonify({"success": False, "message": "Er is een onverwachte fout opgetreden"}), 500

@app.route('/logboek')
def logboek():
    try:
        is_admin = 'is_admin' in session and session['is_admin']
        header = generate_header(is_admin)
        
        months = [
            (9, 2024), (10, 2024), (11, 2024), (12, 2024), (1, 2025)
        ]

        logboek_content = ''
        for i, (month, year) in enumerate(months):
            month_name = calendar.month_name[month]
            content = read_content('', month=month, year=year)
            if not content:
                if month == 9 and year == 2024:
                    content = '''
                    <h3>Week 1</h3>
                    <p>Introductie minor, kennismaking met medestudenten en docenten. Start met het opzetten van de portfolio website.</p>
                    
                    <h3>Week 2</h3>
                    <p>Verdieping in webontwikkeling technieken. Verdere uitwerking van de portfolio structuur en design.</p>
                    
                    <h3>Week 3</h3>
                    <p>Implementatie van basis functionaliteiten op de website, zoals navigatie en contentbeheer. Begin van het verzamelen van projectmateriaal.</p>
                    '''
                else:
                    content = f'<p>Geen logboekinhoud beschikbaar voor {month_name} {year}.</p>'
            
            logboek_content += f'''
            <div class="logboek-month" id="month-{year}-{month}">
                <h2>{month_name} {year}</h2>
                <div class="logboek-entry editable-content">
                    {content}
                </div>
            </div>
            '''

        return f'''
        <html>
            <head>
                <title>Logboek - Portfolio Wessel</title>
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/5.15.3/css/all.min.css">
                <style>
                    {common_styles}
                    body, html {{
                        margin: 0;
                        padding: 0;
                        overflow: hidden;
                    }}
                    .logboek-container {{
                        height: 100vh;
                        overflow-y: scroll;
                        scroll-snap-type: y mandatory;
                    }}
                    .logboek-month {{
                        height: 100vh;
                        scroll-snap-align: start;
                        display: flex;
                        flex-direction: column;
                        padding: 20px;
                        box-sizing: border-box;
                        background-color: var(--bg-color);
                        transition: background-color 0.3s ease;
                    }}
                    .logboek-month h2 {{
                        color: var(--text-color);
                        margin-top: 60px;
                    }}
                    .logboek-entry {{
                        flex-grow: 1;
                        overflow-y: auto;
                        padding: 20px;
                        background-color: var(--content-bg);
                        color: var(--text-color);
                        border-radius: 10px;
                        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                        transition: background-color 0.3s ease, color 0.3s ease;
                    }}
                    .logboek-entry p {{
                        color: var(--text-color);
                        line-height: 1.5;
                    }}
                    .month-indicator {{
                        position: fixed;
                        top: 20px;
                        left: 50%;
                        transform: translateX(-50%);
                        background-color: var(--header-bg);
                        color: var(--text-color);
                        padding: 5px 10px;
                        border-radius: 10px;
                        font-size: 14px;
                        z-index: 1000;
                    }}
                    .editable-content {{
                        white-space: pre-wrap;
                    }}
                    /* Dark mode specifieke stijlen */
                    .dark-mode {{
                        --bg-color: #1d1d1f;
                        --text-color: #f5f5f7;
                        --header-bg: #2c2c2e;
                        --content-bg: #2c2c2e;
                    }}
                    .dark-mode .logboek-month {{
                        background-color: var(--bg-color);
                    }}
                    .dark-mode .logboek-entry {{
                        background-color: var(--content-bg);
                        color: var(--text-color);
                    }}
                    .dark-mode .month-indicator {{
                        background-color: var(--header-bg);
                        color: var(--text-color);
                    }}
                    .dark-mode h1, .dark-mode h2, .dark-mode h3, .dark-mode h4, .dark-mode h5, .dark-mode h6 {{
                        color: var(--text-color);
                    }}
                </style>
            </head>
            <body>
                {header}
                <div class="month-indicator"></div>
                <div class="logboek-container">
                    {logboek_content}
                </div>
                {generate_edit_button(is_admin)}
                {generate_edit_script()}
                {generate_dark_mode_script()}
                <script>
                    const logboekContainer = document.querySelector('.logboek-container');
                    const months = Array.from(document.querySelectorAll('.logboek-month'));
                    const monthIndicator = document.querySelector('.month-indicator');

                    function updateMonthIndicator() {{
                        const scrollPosition = logboekContainer.scrollTop;
                        const pageHeight = window.innerHeight;
                        
                        const currentMonthIndex = Math.floor(scrollPosition / pageHeight);
                        const currentMonth = months[currentMonthIndex];
                        
                        if (currentMonth) {{
                            const monthName = currentMonth.querySelector('h2').textContent;
                            monthIndicator.textContent = monthName;
                        }}
                    }}

                    logboekContainer.addEventListener('scroll', updateMonthIndicator);

                    function scrollToMonth(direction) {{
                        const currentScrollPosition = logboekContainer.scrollTop;
                        const pageHeight = window.innerHeight;
                        const targetScrollPosition = direction === 'next'
                            ? Math.ceil(currentScrollPosition / pageHeight) * pageHeight
                            : Math.floor(currentScrollPosition / pageHeight - 1) * pageHeight;
                        
                        logboekContainer.scrollTo({{
                            top: targetScrollPosition,
                            behavior: 'smooth'
                        }});
                    }}

                    document.addEventListener('keydown', (e) => {{
                        if (e.key === 'ArrowDown' || e.key === 'ArrowRight') {{
                            scrollToMonth('next');
                        }} else if (e.key === 'ArrowUp' || e.key === 'ArrowLeft') {{
                            scrollToMonth('previous');
                        }}
                    }});

                    updateMonthIndicator();

                    function initDarkMode() {{
                        const darkModeEnabled = localStorage.getItem('darkMode') === 'true';
                        document.body.classList.toggle('dark-mode', darkModeEnabled);
                        updateLogoFilter();
                    }}

                    document.addEventListener('DOMContentLoaded', () => {{
                        initDarkMode();
                    }});
                </script>
            </body>
        </html>
        '''
    except Exception as e:
        logging.error(f"Error in logboek route: {str(e)}")
        return f"Er is een fout opgetreden: {str(e)}", 500

@app.route('/contact')
def contact():
    try:
        is_admin = 'is_admin' in session and session['is_admin']
        header = generate_header(is_admin)
        
        content = read_content('contact.txt')
        if not content:
            content = '<h1>Contact</h1><p>Hier komt de contactinformatie.</p>'

        return f'''
        <html>
            <head>
                <title>Contact - Portfolio Wessel</title>
                <style>{common_styles}</style>
            </head>
            <body>
                {header}
                <div class="container">
                    <div class="content-container">
                        <div class="editable-content">
                            {content}
                        </div>
                    </div>
                </div>
                {generate_edit_button(is_admin)}
                {generate_edit_script()}
                {generate_dark_mode_script()}
            </body>
        </html>
        '''
    except Exception as e:
        logging.error(f"Error in contact route: {str(e)}")
        return jsonify({"success": False, "message": f"Er is een onverwachte fout opgetreden: {str(e)}"}), 500

@app.route('/update_folder_order', methods=['POST'])
@admin_required
def update_folder_order():
    if not verify_csrf_token():
        return jsonify({"success": False, "message": "Ongeldige CSRF-token"}), 403
    
    folder_order = request.json.get('folder_order', [])
    if folder_order:
        try:
            # Hier kun je de logica implementeren om de volgorde van mappen bij te werken
            # Dit kan bijvoorbeeld het bijwerken van een database zijn, of het herschrijven van een configuratiebestand
            # Voor nu geven we alleen een succesbericht terug
            return jsonify({"success": True, "message": "Mapvolgorde succesvol bijgewerkt"})
        except Exception as e:
            return jsonify({"success": False, "message": f"Fout bij het bijwerken van de mapvolgorde: {str(e)}"}), 500
    else:
        return jsonify({"success": False, "message": "Geen mapvolgorde ontvangen"}), 400

if __name__ == '__main__':
    logging.debug("Initializing file structure")
    update_file_structure()  # Initialize file structure when starting the app
    app.run(debug=True)

